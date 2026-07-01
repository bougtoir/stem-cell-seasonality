#!/usr/bin/env python3
"""
Generate graphical abstract for Human Reproduction Opinion submission.

HR requirements:
- Distinct single-panel image
- Short legend/title ~25 words
- Uploaded as separate file

This script creates both PNG and PPTX versions.
"""

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def create_graphical_abstract():
    fig, ax = plt.subplots(1, 1, figsize=(10, 7))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis('off')

    # Title
    ax.text(5, 6.6, 'The Invisible Variables', fontsize=20, fontweight='bold',
            ha='center', va='top', color='#1a472a')
    ax.text(5, 6.15, 'Uncontrolled Environmental Factors in\nStem Cell Differentiation',
            fontsize=12, ha='center', va='top', color='#333333')

    # Left column: Controlled vs Uncontrolled
    controlled_y = 5.2
    ax.text(2.2, controlled_y + 0.25, 'Controlled', fontsize=11, fontweight='bold',
            ha='center', color='#2d7d46')

    controlled_vars = ['Temperature 37\u00b0C', 'CO\u2082 5%', 'O\u2082 (some labs)']
    for i, var in enumerate(controlled_vars):
        y = controlled_y - 0.35 * i - 0.15
        rect = FancyBboxPatch((0.5, y - 0.12), 3.4, 0.28,
                               boxstyle="round,pad=0.05",
                               facecolor='#d4edda', edgecolor='#2d7d46', linewidth=1)
        ax.add_patch(rect)
        ax.text(2.2, y + 0.02, var, fontsize=9, ha='center', va='center', color='#1a472a')

    # Dashed separator
    ax.plot([0.5, 3.9], [3.95, 3.95], '--', color='#4a90d9', linewidth=1.5, alpha=0.7)
    ax.text(2.2, 3.75, 'Uncontrolled', fontsize=11, fontweight='bold',
            ha='center', color='#c0392b')

    uncontrolled_vars = ['Humidity', 'VOCs', 'Light/Photoperiod',
                          'EMF (50/60 Hz)', 'Barometric pressure', 'Vibration']
    for i, var in enumerate(uncontrolled_vars):
        y = 3.5 - 0.35 * i
        rect = FancyBboxPatch((0.5, y - 0.12), 3.4, 0.28,
                               boxstyle="round,pad=0.05",
                               facecolor='#f8d7da', edgecolor='#c0392b', linewidth=1)
        ax.add_patch(rect)
        ax.text(2.2, y + 0.02, var, fontsize=9, ha='center', va='center', color='#721c24')

    # Arrow
    ax.annotate('', xy=(5.3, 4.0), xytext=(4.2, 4.0),
                arrowprops=dict(arrowstyle='->', color='#555555', lw=2))

    # Right column: Key Findings
    ax.text(7.5, 5.5, 'Key Findings', fontsize=13, fontweight='bold',
            ha='center', color='#333333')

    findings = [
        ('\u2460 PSCs are uniquely vulnerable',
         'Several-fold more sensitive than\nsomatic cells to environmental stress',
         '#fff3cd', '#856404'),
        ('\u2461 "Seasonality" = fiscal year artifact',
         'GEO March peak driven by Japan\'s\nfiscal year-end, not biology',
         '#d6e9f8', '#1a5276'),
        ('\u2462 Prospective monitoring needed',
         'IoT sensors + outcome tracking\nfor \u226512 months across latitudes',
         '#d5f5e3', '#1a7a3a'),
    ]

    for i, (title, desc, bg, tc) in enumerate(findings):
        y_center = 4.9 - 1.15 * i
        rect = FancyBboxPatch((5.3, y_center - 0.4), 4.2, 0.85,
                               boxstyle="round,pad=0.1",
                               facecolor=bg, edgecolor=tc, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(7.4, y_center + 0.15, title, fontsize=9.5, fontweight='bold',
                ha='center', va='center', color=tc)
        ax.text(7.4, y_center - 0.15, desc, fontsize=8, ha='center', va='center',
                color='#444444', linespacing=1.3)

    # Bottom quote
    rect = FancyBboxPatch((0.5, 0.3), 9.0, 0.8,
                           boxstyle="round,pad=0.1",
                           facecolor='#f0e6f6', edgecolor='#6c3483', linewidth=1.5)
    ax.add_patch(rect)
    ax.text(5, 0.85, 'We spent two decades optimizing recipes', fontsize=10,
            ha='center', va='center', color='#4a235a', style='italic')
    ax.text(5, 0.55, 'while ignoring the kitchen.', fontsize=11,
            ha='center', va='center', color='#6c3483', fontweight='bold', style='italic')

    plt.tight_layout()

    # Save PNG
    png_path = os.path.join(OUTPUT_DIR, "HumReprod_Graphical_Abstract.png")
    fig.savefig(png_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Graphical abstract PNG saved to: {png_path}")

    # Create PPTX
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                      PptxInches(12.3), PptxInches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Graphical Abstract"
    p.font.size = PptxPt(24)
    p.font.bold = True

    # Image
    slide.shapes.add_picture(png_path, PptxInches(1.5), PptxInches(0.8),
                              PptxInches(10.3), PptxInches(5.8))

    # Caption
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.8),
                                      PptxInches(12.3), PptxInches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ("Graphical Abstract. Stem cell laboratories control temperature and CO\u2082 "
              "but leave multiple environmental variables unmonitored; apparent database "
              "seasonality is an institutional artifact.")
    p.font.size = PptxPt(11)

    pptx_path = os.path.join(OUTPUT_DIR, "HumReprod_Graphical_Abstract.pptx")
    prs.save(pptx_path)
    print(f"Graphical abstract PPTX saved to: {pptx_path}")

    return png_path, pptx_path


if __name__ == "__main__":
    create_graphical_abstract()
