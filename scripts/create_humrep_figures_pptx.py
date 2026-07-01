#!/usr/bin/env python3
"""
Generate figures PPTX for Human Reproduction Opinion submission.

HR requirements:
- Figures uploaded as separate files
- Figure legends in manuscript text (not in figure files)
- Each figure self-contained

Creates one PPTX with Figure 1 and Figure 2 on separate slides.
Also creates individual PNG files for upload.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def create_figures_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Figure 1: Natural experiment
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 1. Academic calendar variation as a natural experiment"
    p.font.size = Pt(20)
    p.font.bold = True

    fig1_path = os.path.join(OUTPUT_DIR, "natural_experiment_figure.png")
    if os.path.exists(fig1_path):
        slide1.shapes.add_picture(fig1_path, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(6.3))
    else:
        txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(3),
                                            Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "[Natural experiment figure - run analysis/natural_experiment.py first]"
        p2.font.size = Pt(16)
        p2.alignment = PP_ALIGN.CENTER

    # Figure 2: Research roadmap (reuse GEO solar comprehensive as proxy,
    # or create placeholder)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 2. Research roadmap for environmental profiling of PSC culture"
    p.font.size = Pt(20)
    p.font.bold = True

    # Use the existing CellStemCell figures PPTX figure 2 if available,
    # or the comprehensive GEO solar figure
    fig2_path = os.path.join(OUTPUT_DIR, "geo_solar_comprehensive.png")
    if os.path.exists(fig2_path):
        slide2.shapes.add_picture(fig2_path, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(6.3))
    else:
        txBox2 = slide2.shapes.add_textbox(Inches(2), Inches(3),
                                            Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "[Research roadmap figure - to be created]"
        p2.font.size = Pt(16)
        p2.alignment = PP_ALIGN.CENTER

    pptx_path = os.path.join(OUTPUT_DIR, "HumReprod_Figures.pptx")
    prs.save(pptx_path)
    print(f"Figures PPTX saved to: {pptx_path}")

    # Copy individual figure PNGs with HR-specific names
    import shutil
    if os.path.exists(fig1_path):
        fig1_dest = os.path.join(OUTPUT_DIR, "HumReprod_Figure1.png")
        shutil.copy2(fig1_path, fig1_dest)
        print(f"Figure 1 PNG: {fig1_dest}")
    if os.path.exists(fig2_path):
        fig2_dest = os.path.join(OUTPUT_DIR, "HumReprod_Figure2.png")
        shutil.copy2(fig2_path, fig2_dest)
        print(f"Figure 2 PNG: {fig2_dest}")

    return pptx_path


if __name__ == "__main__":
    create_figures_pptx()
