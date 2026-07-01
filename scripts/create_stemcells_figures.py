#!/usr/bin/env python3
"""
Generate figures for STEM CELLS Perspective submission.

Creates:
- StemCells_Figure1.png/tiff  (GEO natural experiment)
- StemCells_Figure2.png/tiff  (research roadmap)
- StemCells_Figures.pptx      (editable, one figure per slide)

Per journal guidelines, figures are submitted as separate files,
not embedded in the manuscript.
"""

import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

FIG1_SRC = os.path.join(OUTPUT_DIR, "natural_experiment_figure.png")
FIG2_SRC = os.path.join(OUTPUT_DIR, "geo_solar_comprehensive.png")


def png_to_tiff(png_path, tiff_path):
    """Convert PNG to TIFF (LZW compression, 300 dpi)."""
    img = Image.open(png_path)
    if img.mode == 'RGBA':
        bg = Image.new('RGB', img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[3])
        img = bg
    elif img.mode != 'RGB':
        img = img.convert('RGB')
    img.save(tiff_path, format='TIFF', compression='tiff_lzw', dpi=(300, 300))
    print(f"  TIFF: {tiff_path}")


def create_figures():
    # ── Figure 1: GEO natural experiment ──
    fig1_png = os.path.join(OUTPUT_DIR, "StemCells_Figure1.png")
    fig1_tiff = os.path.join(OUTPUT_DIR, "StemCells_Figure1.tiff")

    if os.path.exists(FIG1_SRC):
        shutil.copy2(FIG1_SRC, fig1_png)
        print(f"Figure 1 PNG: {fig1_png}")
        png_to_tiff(fig1_png, fig1_tiff)
    else:
        print(f"WARNING: Source figure not found: {FIG1_SRC}")

    # ── Figure 2: Research roadmap ──
    fig2_png = os.path.join(OUTPUT_DIR, "StemCells_Figure2.png")
    fig2_tiff = os.path.join(OUTPUT_DIR, "StemCells_Figure2.tiff")

    if os.path.exists(FIG2_SRC):
        shutil.copy2(FIG2_SRC, fig2_png)
        print(f"Figure 2 PNG: {fig2_png}")
        png_to_tiff(fig2_png, fig2_tiff)
    else:
        print(f"WARNING: Source figure not found: {FIG2_SRC}")

    # ── Editable PPTX ──
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Figure 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 1. Institutional calendar confounding in GEO PSC datasets"
    p.font.size = Pt(20)
    p.font.bold = True

    if os.path.exists(fig1_png):
        slide1.shapes.add_picture(fig1_png, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(5.5))

    # Caption
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(6.6),
                                       Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "(A) Monthly submission distributions for 6,101 GEO PSC differentiation "
        "datasets grouped by country academic-year start. (B) Country-level heatmap "
        "of submissions. (C) Circular Rayleigh vectors showing mean submission "
        "direction. (D) Statistical summary: only Japan/Korea shows significant "
        "seasonality, aligned with Japan's March 31 fiscal year-end."
    )
    p.font.size = Pt(10)

    # Slide 2: Figure 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 2. Proposed research roadmap for environmental profiling"
    p.font.size = Pt(20)
    p.font.bold = True

    if os.path.exists(fig2_png):
        slide2.shapes.add_picture(fig2_png, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(5.5))

    # Caption
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(6.6),
                                       Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "(A) Evidence matrix classifying uncontrolled environmental variables by "
        "hemisphere-dependence and strength of existing evidence. (B) Three-phase "
        "investigation strategy: Phase I (passive IoT monitoring), Phase II "
        "(IVF registry analysis), Phase III (controlled intervention)."
    )
    p.font.size = Pt(10)

    pptx_path = os.path.join(OUTPUT_DIR, "StemCells_Figures.pptx")
    prs.save(pptx_path)
    print(f"Figures PPTX: {pptx_path}")

    return fig1_png, fig2_png, pptx_path


if __name__ == "__main__":
    create_figures()
