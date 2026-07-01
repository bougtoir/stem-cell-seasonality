#!/usr/bin/env python3
"""
Generate figures for Nature Methods Perspective submission.

Creates:
- NatMethods_Figure1.png/tiff  (natural experiment)
- NatMethods_Figure2.png/tiff  (research roadmap)
- NatMethods_Figures.pptx      (editable, one figure per slide)

Source figures come from the analysis/ output directory.
"""

import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

FIG1_SRC = os.path.join(OUTPUT_DIR, "natural_experiment_figure.png")
FIG2_SRC = os.path.join(OUTPUT_DIR, "geo_solar_comprehensive.png")


def png_to_tiff(png_path, tiff_path):
    """Convert PNG to TIFF (LZW compression, 300 dpi)."""
    img = Image.open(png_path)
    if img.mode == 'RGBA':
        bg = Image.new('RGB', img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[3])
        img = bg
    elif img.mode != 'RGB':
        img = img.convert('RGB')
    img.save(tiff_path, format='TIFF', compression='tiff_lzw', dpi=(300, 300))
    print(f"  TIFF: {tiff_path}")


def create_figures():
    # ── Figure 1: Natural experiment ──
    fig1_png = os.path.join(OUTPUT_DIR, "NatMethods_Figure1.png")
    fig1_tiff = os.path.join(OUTPUT_DIR, "NatMethods_Figure1.tiff")

    if os.path.exists(FIG1_SRC):
        shutil.copy2(FIG1_SRC, fig1_png)
        print(f"Figure 1 PNG: {fig1_png}")
        png_to_tiff(fig1_png, fig1_tiff)
    else:
        print(f"WARNING: Source figure not found: {FIG1_SRC}")

    # ── Figure 2: Research roadmap ──
    fig2_png = os.path.join(OUTPUT_DIR, "NatMethods_Figure2.png")
    fig2_tiff = os.path.join(OUTPUT_DIR, "NatMethods_Figure2.tiff")

    if os.path.exists(FIG2_SRC):
        shutil.copy2(FIG2_SRC, fig2_png)
        print(f"Figure 2 PNG: {fig2_png}")
        png_to_tiff(fig2_png, fig2_tiff)
    else:
        print(f"WARNING: Source figure not found: {FIG2_SRC}")

    # ── Editable PPTX ──
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Figure 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 1 | Academic calendar variation as a natural experiment"
    p.font.size = Pt(20)
    p.font.bold = True

    if os.path.exists(fig1_png):
        slide1.shapes.add_picture(fig1_png, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(5.5))

    # Caption
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(6.6),
                                       Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("a, Normalized monthly distributions of GEO PSC dataset submissions by "
              "academic year group. b, Country-level heatmap. c, Circular mean directions "
              "(Rayleigh vectors). d, Statistical summary. n=6,101 datasets; 3,195 with "
              "country affiliation.")
    p.font.size = Pt(10)

    # Slide 2: Figure 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                       Inches(12.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 2 | Research roadmap for environmental profiling of PSC culture"
    p.font.size = Pt(20)
    p.font.bold = True

    if os.path.exists(fig2_png):
        slide2.shapes.add_picture(fig2_png, Inches(0.5), Inches(0.9),
                                   Inches(12.3), Inches(5.5))

    # Caption
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(6.6),
                                       Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("a, Evidence matrix: uncontrolled environmental variables classified by "
              "hemisphere-dependence and evidence level. b, Three-phase investigation "
              "strategy from passive monitoring through data mining to controlled "
              "intervention.")
    p.font.size = Pt(10)

    pptx_path = os.path.join(OUTPUT_DIR, "NatMethods_Figures.pptx")
    prs.save(pptx_path)
    print(f"Figures PPTX: {pptx_path}")

    return fig1_png, fig2_png, pptx_path


if __name__ == "__main__":
    create_figures()
