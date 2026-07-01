#!/usr/bin/env python3
"""
Generate editable PPTX figures for the GEO × Solar Activity analysis paper.

Figure 1: 6-panel comprehensive analysis (from PNG)
Figure 2: Study design schematic (limitations & future directions)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(SCRIPT_DIR)
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, italic=False, color=RGBColor(0, 0, 0),
                alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=10, font_color=RGBColor(255, 255, 255), bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=10, font_color=RGBColor(0, 0, 0), bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    return shape


def create_figure1(prs):
    """Figure 1: Comprehensive 6-panel analysis figure (embed PNG)."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                "Figure 1. Temporal and Geographic Patterns in GEO PSC Datasets",
                font_size=18, bold=True)

    # Embed the comprehensive figure
    fig_path = os.path.join(OUTPUT_DIR, "geo_solar_comprehensive.png")
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path,
                                 Inches(0.5), Inches(0.7),
                                 width=Inches(12.3))

    # Caption
    add_textbox(slide, Inches(0.3), Inches(6.9), Inches(12.5), Inches(0.6),
                "(A) Monthly distribution with chi-square test. (B) Country distribution. "
                "(C) Hemisphere comparison. (D) Annual trend vs solar cycle. "
                "(E) F10.7 correlation. (F) Detrended SSN correlation.",
                font_size=9, italic=True)


def create_figure2(prs):
    """Figure 2: Study design — data flow, limitations, and next steps."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                "Figure 2. Study Design, Limitations, and Proposed Next Steps",
                font_size=18, bold=True)

    # ── Left panel: Data sources ──
    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(3.5), Inches(0.4),
                "A. Data Sources", font_size=14, bold=True,
                color=RGBColor(0, 80, 160))

    sources = [
        ("GEO (NCBI)", "6,101 PSC datasets\n2001\u20132026", RGBColor(200, 225, 255)),
        ("NOAA SWPC", "Solar indices\n(F10.7, SSN)", RGBColor(255, 230, 200)),
        ("GEO SOFT", "Country metadata\n(n=600 sample)", RGBColor(220, 240, 220)),
    ]
    for i, (title, desc, col) in enumerate(sources):
        y = Inches(1.3 + i * 1.3)
        add_rounded_rect(slide, Inches(0.7), y, Inches(3.2), Inches(1.1),
                        col, f"{title}\n{desc}", font_size=11, bold=False)

    # ── Middle panel: Key findings ──
    add_textbox(slide, Inches(4.5), Inches(0.8), Inches(4.0), Inches(0.4),
                "B. Key Findings", font_size=14, bold=True,
                color=RGBColor(0, 120, 0))

    findings = [
        ("\u2713 Seasonal non-uniformity", "\u03c7\u00b2=32.3, p=0.0007",
         RGBColor(200, 240, 200)),
        ("\u2717 Hemisphere comparison", "97% NH, only 3% SH",
         RGBColor(255, 220, 220)),
        ("\u2717 Solar correlation", "Detrended r=0.02, p=0.75",
         RGBColor(255, 220, 220)),
    ]
    for i, (title, detail, col) in enumerate(findings):
        y = Inches(1.3 + i * 1.3)
        add_rounded_rect(slide, Inches(4.7), y, Inches(3.8), Inches(1.1),
                        col, f"{title}\n{detail}", font_size=11, bold=False)

    # ── Right panel: Limitations & next steps ──
    add_textbox(slide, Inches(9.0), Inches(0.8), Inches(4.0), Inches(0.4),
                "C. Limitations & Next Steps", font_size=14, bold=True,
                color=RGBColor(160, 0, 0))

    limitations = [
        ("Submission date \u2260 Experiment date", "6\u201318 month lag obscures signal",
         RGBColor(255, 240, 220)),
        ("Hemisphere imbalance (97:3)", "Cannot test photoperiod inversion",
         RGBColor(255, 240, 220)),
        ("Institutional cycles", "Academic calendar confound",
         RGBColor(255, 240, 220)),
    ]
    for i, (title, detail, col) in enumerate(limitations):
        y = Inches(1.3 + i * 1.0)
        add_rounded_rect(slide, Inches(9.2), y, Inches(3.8), Inches(0.85),
                        col, f"{title}\n{detail}", font_size=10, bold=False)

    # Arrow to next steps
    add_textbox(slide, Inches(9.5), Inches(4.3), Inches(3.5), Inches(0.3),
                "\u2193 Proposed solutions \u2193", font_size=11, bold=True,
                color=RGBColor(100, 0, 0), alignment=PP_ALIGN.CENTER)

    next_steps = [
        ("HFEA Registry (UK)", "Cycle-level IVF data, 51\u00b0N"),
        ("ANZARD (Australia)", "Southern Hemisphere, \u221233\u00b0S"),
        ("Environmental monitoring", "Humidity, light, EMF in PSC labs"),
    ]
    for i, (title, detail) in enumerate(next_steps):
        y = Inches(4.7 + i * 0.85)
        add_rounded_rect(slide, Inches(9.2), y, Inches(3.8), Inches(0.75),
                        RGBColor(220, 235, 255),
                        f"{title}\n{detail}", font_size=10, bold=False)

    # Caption
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(12.5), Inches(0.5),
                "Figure 2. Schematic overview of the study design, key findings, "
                "principal limitations, and proposed next steps for hypothesis testing.",
                font_size=9, italic=True)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_figure1(prs)
    create_figure2(prs)

    out_path = os.path.join(OUTPUT_DIR, "GEO_Solar_Analysis_Figures.pptx")
    prs.save(out_path)
    print(f"PPTX saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    main()
